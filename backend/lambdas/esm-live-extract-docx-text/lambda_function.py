#!/usr/bin/env python
"""Extract plain text from Office documents stored in the global document store.

Handles: .docx (python-docx), .pptx (python-pptx), .txt/.md (UTF-8 decode).
Note: .doc (Word 97-2003) is not supported — requires the antiword system binary
which is unavailable in the Lambda base image.

No OCR required — all supported formats contain embedded text.

Idempotent: if text.txt already exists in the global store, extraction
is skipped and the existing result is returned.

Input (JSON body):
    md5  (str, required): 32-character hex MD5 hash of the document.

Output (200 — success):
    { "md5": "<hash>", "status": "ready", "char_count": 12043 }

Output (200 — already extracted):
    { "md5": "<hash>", "status": "ready", "skipped": true }

Output (400):  missing md5.
Output (404):  document not found in global store.
Output (415):  unsupported file extension.
Output (422):  extraction produced no text.
Output (500):  storage or extraction failure.

Related: philkomarny/SkillStore#15, philkomarny/SkillStore#14
"""

import io
from json import dumps, loads
import os
from datetime import datetime, timezone
from logging import Logger
from typing import Any

import boto3
from botocore.exceptions import ClientError
from docx import Document
from pptx import Presentation

from skillstore_base import configure_logger, isnullstr

BUCKET_NAME: str = os.getenv("BUCKET_NAME", "mskillsiq")

# S3 layout — https://github.com/philkomarny/SkillStore/issues/16:
#   store      — raw uploads:  eduskillsmp/documents/store/<md5[:2]>/<md5>.json
#                               eduskillsmp/documents/store/<md5[:2]>/<md5>.<ext>
#   text-store — extracted:    eduskillsmp/documents/text-store/<md5[:2]>/<md5>.txt
_STORE_PREFIX = "eduskillsmp/documents/store"
_TEXT_STORE_PREFIX = "eduskillsmp/documents/text-store"

_SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".txt", ".md"}

s3_client = boto3.client("s3")
logger: Logger = configure_logger(__name__)


def _metadata_key(md5: str) -> str:
    """Return the S3 key for a document's metadata.json."""
    return f"{_STORE_PREFIX}/{md5[:2]}/{md5}.json"


def _original_key(md5: str, ext: str) -> str:
    """Return the S3 key for a document's original file bytes."""
    return f"{_STORE_PREFIX}/{md5[:2]}/{md5}{ext}"


def _text_key(md5: str) -> str:
    """Return the S3 key for a document's extracted text."""
    return f"{_TEXT_STORE_PREFIX}/{md5[:2]}/{md5}.txt"


def _text_already_extracted(md5: str) -> bool:
    """Return True if text.txt already exists — idempotent retry guard."""
    try:
        s3_client.head_object(Bucket=BUCKET_NAME, Key=_text_key(md5))
        return True
    except ClientError as exc:
        if exc.response["Error"]["Code"] == "404":
            return False
        raise


def _read_metadata(md5: str) -> dict:
    """Read and return the document's metadata.json."""
    obj = s3_client.get_object(Bucket=BUCKET_NAME, Key=_metadata_key(md5))
    return loads(obj["Body"].read())


def _read_original(md5: str, ext: str) -> bytes:
    """Read the original file bytes from S3."""
    obj = s3_client.get_object(Bucket=BUCKET_NAME, Key=_original_key(md5, ext))
    return obj["Body"].read()


def _write_text(md5: str, text: str) -> None:
    """Write extracted plain text to text.txt in the global store."""
    s3_client.put_object(
        Bucket=BUCKET_NAME,
        Key=_text_key(md5),
        Body=text.encode("utf-8"),
        ContentType="text/plain",
    )


def _update_metadata(md5: str, metadata: dict, status: str) -> None:
    """Update metadata.json with new status and extraction timestamp."""
    metadata["status"] = status
    metadata["text_extracted"] = status == "ready"
    metadata["extracted_at"] = datetime.now(tz=timezone.utc).isoformat()
    s3_client.put_object(
        Bucket=BUCKET_NAME,
        Key=_metadata_key(md5),
        Body=dumps(metadata, indent=2),
        ContentType="application/json",
    )


def _extract_docx(raw_bytes: bytes) -> str:
    """Extract plain text from a .docx file.

    Iterates paragraphs and table cells to capture all text content.
    """
    doc = Document(io.BytesIO(raw_bytes))
    parts = []

    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text.strip())

    return "\n\n".join(parts)


def _extract_pptx(raw_bytes: bytes) -> str:
    """Extract plain text from a .pptx file.

    Iterates all slides and text frames.
    """
    prs = Presentation(io.BytesIO(raw_bytes))
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    slide_parts.append(text)
        if slide_parts:
            slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(slide_parts))

    return "\n\n".join(slides_text)


def _extract_text_file(raw_bytes: bytes) -> str:
    """Decode a plain text or markdown file as UTF-8."""
    return raw_bytes.decode("utf-8", errors="replace")


def _error(message: str, status: int = 500) -> dict[str, Any]:
    logger.error(message)
    return {"statusCode": status, "body": dumps({"message": message})}


def handler(event: dict[str, Any], _) -> dict[str, Any]:
    """Extract plain text from a DOCX, PPTX, TXT, or MD document.

    Reads the original file from the global store, routes to the
    appropriate extractor based on file extension, writes text.txt,
    and flips metadata.json status to 'ready'.

    Args:
        event: Lambda invocation event with JSON body fields.

    Returns:
        200 { md5, status, char_count } on success.
        200 { md5, status, skipped } if already extracted.
        400 { message } for missing md5.
        404 { message } when document not found.
        415 { message } for unsupported extension.
        422 { message } when extraction yields no text.
        500 { message } on failure.
    """
    logger.info(f"Incoming Event: {dumps(event)}")

    params = event
    if isinstance(event.get("body"), str):
        try:
            params = loads(event["body"])
        except Exception as exc:
            return _error(f"Invalid JSON body: {exc}", status=400)

    md5: str = (params.get("md5") or "").strip().lower()
    if isnullstr(md5):
        return _error("Missing required parameter: 'md5'", status=400)

    logger.info(f"Resolved Params: {dumps({'md5': md5})}")

    # Idempotency: skip if text.txt already exists.
    if _text_already_extracted(md5):
        logger.info(f"text.txt already exists for md5={md5}; skipping")
        try:
            metadata = _read_metadata(md5)
            if metadata.get("status") != "ready":
                _update_metadata(md5, metadata, "ready")
        except Exception:
            pass
        return {"statusCode": 200, "body": dumps({"md5": md5, "status": "ready", "skipped": True})}

    try:
        metadata = _read_metadata(md5)
    except ClientError as exc:
        if exc.response["Error"]["Code"] in ("404", "NoSuchKey"):
            return _error(f"Document not found: {md5}", status=404)
        raise

    ext = metadata.get("ext", "").lower()
    if ext not in _SUPPORTED_EXTENSIONS:
        return _error(f"Unsupported file type: '{ext}' — supported: {sorted(_SUPPORTED_EXTENSIONS)}", status=415)

    filename = metadata.get("filename", "unknown")
    size_bytes = metadata.get("size_bytes", 0)
    logger.info(f"Extracting office document: md5={md5} filename={filename} ext={ext} size={size_bytes}b")

    try:
        raw_bytes = _read_original(md5, ext)
    except ClientError as exc:
        if exc.response["Error"]["Code"] in ("404", "NoSuchKey"):
            return _error(f"Original file not found for md5={md5}", status=404)
        raise

    # Route to the appropriate extractor.
    if ext == ".docx":
        text = _extract_docx(raw_bytes)
    elif ext == ".pptx":
        text = _extract_pptx(raw_bytes)
    else:  # .txt, .md
        text = _extract_text_file(raw_bytes)

    if not text.strip():
        _update_metadata(md5, metadata, "error")
        return _error(f"Extraction produced no text for md5={md5}", status=422)

    # Write order: text first, then status flip — idempotent on retry.
    _write_text(md5, text)
    _update_metadata(md5, metadata, "ready")

    logger.info(f"Office extraction complete: md5={md5} ext={ext} chars={len(text)}")
    result = {"statusCode": 200, "body": dumps({"md5": md5, "status": "ready", "char_count": len(text)})}
    logger.info(f"Return Value: {dumps(result)}")
    return result
